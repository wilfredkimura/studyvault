import os
import sys
import hashlib
import re

# Third-party libraries
try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    import pypdfium2 as pdfium
except ImportError:
    pdfium = None

try:
    import pytesseract
except ImportError:
    pytesseract = None

try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None


def compute_sha256(file_path: str) -> str:
    """Compute the SHA-256 hash of a file's binary contents."""
    sha256_hash = hashlib.sha256()
    with open(file_path, "rb") as f:
        for byte_block in iter(lambda: f.read(4096), b""):
            sha256_hash.update(byte_block)
    return sha256_hash.hexdigest()


def extract_txt(file_path: str) -> str:
    """Extract plain text from txt or md files."""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def extract_docx(file_path: str) -> str:
    """Extract paragraph and table text from .docx files."""
    if docx is None:
        return "[Error: python-docx not installed]"
    try:
        doc = docx.Document(file_path)
        text_list = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_list.append(paragraph.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    text_list.append(" | ".join(row_text))
        return "\n".join(text_list)
    except Exception as e:
        return f"[Error parsing .docx: {str(e)}]"


def extract_pptx(file_path: str) -> str:
    """Extract slide text from .pptx files."""
    if Presentation is None:
        return "[Error: python-pptx not installed]"
    try:
        prs = Presentation(file_path)
        text_list = []
        for slide_idx, slide in enumerate(prs.slides):
            text_list.append(f"--- Slide {slide_idx + 1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_list.append(shape.text)
        return "\n".join(text_list)
    except Exception as e:
        return f"[Error parsing .pptx: {str(e)}]"


def extract_xlsx(file_path: str) -> str:
    """Extract sheet row text from .xlsx files."""
    if openpyxl is None:
        return "[Error: openpyxl not installed]"
    try:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        text_list = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_list.append(f"--- Sheet: {sheet_name} ---")
            for row in sheet.iter_rows(values_only=True):
                row_vals = [str(cell).strip() for cell in row if cell is not None]
                if row_vals:
                    text_list.append(" | ".join(row_vals))
        return "\n".join(text_list)
    except Exception as e:
        return f"[Error parsing .xlsx: {str(e)}]"


def extract_pdf_ocr(file_path: str) -> str:
    """Run OCR page-by-page using pypdfium2 and pytesseract."""
    if pdfium is None:
        return "[OCR Error: pypdfium2 is not installed]"
    if pytesseract is None:
        return "[OCR Error: pytesseract is not installed]"
        
    ocr_text_list = []
    try:
        pdf = pdfium.PdfDocument(file_path)
        for idx in range(len(pdf)):
            page = pdf[idx]
            # Render at 150 DPI (scale=2.0) for good balance of speed and quality
            bitmap = page.render(scale=2.0)
            pil_image = bitmap.to_pil()
            text = pytesseract.image_to_string(pil_image)
            if text.strip():
                ocr_text_list.append(f"--- Page {idx + 1} (OCR) ---\n{text}")
        pdf.close()
    except Exception as e:
        return f"[OCR Failed: {str(e)}]"
        
    return "\n".join(ocr_text_list)


def extract_pdf(file_path: str) -> str:
    """Extract text from .pdf using pdfplumber, falling back to OCR if empty."""
    if pdfplumber is None:
        return "[Error: pdfplumber not installed]"
        
    text_list = []
    has_text = False
    try:
        with pdfplumber.open(file_path) as pdf:
            for idx, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text and text.strip():
                    text_list.append(f"--- Page {idx + 1} ---\n{text}")
                    has_text = True
    except Exception as e:
        print(f"pdfplumber error: {e}", file=sys.stderr)
        
    if not has_text:
        # Fallback to OCR
        print("No text found, running OCR on PDF...", file=sys.stderr)
        ocr_result = extract_pdf_ocr(file_path)
        if ocr_result and not ocr_result.startswith("[OCR"):
            return ocr_result
            
    return "\n".join(text_list)


def extract_fallback_binary(file_path: str) -> str:
    """Simple 'strings'-like extraction fallback for legacy formats (.doc, .ppt, .xls)."""
    try:
        with open(file_path, "rb") as f:
            data = f.read()
        # Find runs of 4 or more printable ASCII characters
        words = re.findall(br'[a-zA-Z0-9\s\.,;:!@#\$%\^&\*\(\)\-_\+=\[\]\{\}<>\?\/\\|`~]{4,}', data)
        text = "\n".join(w.decode("ascii", errors="ignore").strip() for w in words)
        return f"[Extracted from Legacy Binary Format]\n{text}"
    except Exception as e:
        return f"[Error parsing legacy binary file: {str(e)}]"


def extract_text_and_hash(file_path: str, file_type: str) -> tuple[str, str]:
    """Extract text and compute SHA-256 hash for a given file path and type."""
    if not os.path.exists(file_path):
        return f"[Error: File not found at {file_path}]", "empty_hash"
        
    # 1. Compute hash
    try:
        file_hash = compute_sha256(file_path)
    except Exception as e:
        file_hash = f"hash_error_{hash(file_path)}"
        
    # 2. Extract text based on type
    file_type = file_type.lower()
    
    if file_type in ["txt", "md"]:
        text = extract_txt(file_path)
    elif file_type == "pdf":
        text = extract_pdf(file_path)
    elif file_type == "docx":
        text = extract_docx(file_path)
    elif file_type == "pptx":
        text = extract_pptx(file_path)
    elif file_type == "xlsx":
        text = extract_xlsx(file_path)
    elif file_type in ["doc", "ppt", "xls"]:
        # Try modern parser first (in case it was renamed docx/pptx/xlsx)
        if file_type == "doc":
            text = extract_docx(file_path)
            if text.startswith("[Error"):
                text = extract_fallback_binary(file_path)
        elif file_type == "ppt":
            text = extract_pptx(file_path)
            if text.startswith("[Error"):
                text = extract_fallback_binary(file_path)
        elif file_type == "xls":
            text = extract_xlsx(file_path)
            if text.startswith("[Error"):
                text = extract_fallback_binary(file_path)
    else:
        # Fallback binary reader
        text = extract_fallback_binary(file_path)
        
    return text, file_hash
