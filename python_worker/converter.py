import os
import re
import time
import hashlib
import tempfile
import sys

# Third-party libraries
try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from fpdf import FPDF
except ImportError:
    FPDF = None


class StudyVaultPDF(FPDF):
    def header(self):
        if getattr(self, 'is_slide', False):
            return
        if hasattr(self, 'doc_title'):
            self.set_font('helvetica', 'B', 8)
            self.set_text_color(120, 120, 120)
            self.cell(0, 10, self.doc_title, border=0, align='R')
            self.ln(10)
            
    def footer(self):
        if getattr(self, 'is_slide', False):
            return
        self.set_y(-15)
        self.set_font('helvetica', 'I', 8)
        self.set_text_color(120, 120, 120)
        self.cell(0, 10, f'Page {self.page_no()}/{{nb}}', border=0, align='C')


def clean_pdf_text(text: str) -> str:
    """Sanitize text to fit standard Latin-1 encoding used in FPDF base fonts."""
    if not text:
        return ""
    replacements = {
        '\u2018': "'",  # Smart single quote open
        '\u2019': "'",  # Smart single quote close
        '\u201c': '"',  # Smart double quote open
        '\u201d': '"',  # Smart double quote close
        '\u2013': '-',  # En dash
        '\u2014': '-',  # Em dash
        '\u2022': '*',  # Bullet point
        '\u2026': '...', # Ellipsis
    }
    for orig, repl in replacements.items():
        text = text.replace(orig, repl)
        
    # Encode as latin-1, replacing unknown chars with '?'
    text_bytes = text.encode('latin-1', errors='replace')
    return text_bytes.decode('latin-1')


def clean_extracted_text(text: str) -> str:
    """Remove internal Page headers/footers added by pdfplumber or OCR."""
    if not text:
        return ""
    text = re.sub(r'--- Page \d+ (?:OCR)?---\n', '', text)
    return text


def write_txt(text: str, output_path: str, source_name: str = ""):
    """Write plain text file with conversion header."""
    header = f"[Converted from {source_name}]\n\n" if source_name else ""
    with open(output_path, "w", encoding="utf-8", errors="ignore") as f:
        f.write(header + text)


def write_md(text: str, output_path: str, source_name: str = ""):
    """Write markdown file with conversion header."""
    header = f"# Converted from {source_name}\n\n" if source_name else ""
    with open(output_path, "w", encoding="utf-8", errors="ignore") as f:
        f.write(header + text)


def write_docx(text: str, output_path: str):
    """Write real Word (.docx) document."""
    if docx is None:
        raise ImportError("python-docx is not installed in the python environment.")
        
    doc = docx.Document()
    # Set margins
    for section in doc.sections:
        section.top_margin = docx.shared.Inches(1)
        section.bottom_margin = docx.shared.Inches(1)
        section.left_margin = docx.shared.Inches(1)
        section.right_margin = docx.shared.Inches(1)
        
    paragraphs = text.split('\n')
    for p in paragraphs:
        if p.strip():
            doc.add_paragraph(p)
        else:
            # Add empty line spacing
            doc.add_paragraph("")
            
    doc.save(output_path)


def write_pptx(text: str, output_path: str):
    """Write real PowerPoint (.pptx) presentation."""
    if Presentation is None:
        raise ImportError("python-pptx is not installed in the python environment.")
        
    prs = Presentation()
    blank_layout = prs.slide_layouts[6] # Blank slide layout
    
    # Split text by page/slide markers
    page_regex = r'--- (?:Page \d+(?: \(OCR\))?|Slide \d+|Sheet: [^\n\-]+) ---'
    markers = list(re.finditer(page_regex, text))
    
    slide_contents = []
    if markers:
        for i, marker in enumerate(markers):
            if i == 0:
                pre_text = text[:marker.start()].strip()
                if pre_text:
                    slide_contents.append(pre_text)
            
            start_pos = marker.end()
            end_pos = markers[i+1].start() if i + 1 < len(markers) else len(text)
            content = text[start_pos:end_pos].strip()
            if content:
                slide_contents.append(content)
            else:
                slide_contents.append(f"Slide {i+1}")
    else:
        # Fallback split by paragraphs or chunk sizes
        paragraphs = [p.strip() for p in text.split('\n') if p.strip()]
        current_chunk = []
        current_len = 0
        for p in paragraphs:
            if current_len + len(p) > 600:
                if current_chunk:
                    slide_contents.append("\n\n".join(current_chunk))
                    current_chunk = []
                    current_len = 0
            current_chunk.append(p)
            current_len += len(p)
        if current_chunk:
            slide_contents.append("\n\n".join(current_chunk))
            
    if not slide_contents:
        slide_contents = ["Empty presentation."]
        
    from pptx.util import Inches, Pt
    for content in slide_contents:
        slide = prs.slides.add_slide(blank_layout)
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9.0)
        height = Inches(6.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        paragraphs = content.split('\n')
        for idx, p_text in enumerate(paragraphs):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = p_text
            p.font.size = Pt(14)
            p.font.name = 'Arial'
            
    prs.save(output_path)



def write_pdf(text: str, output_path: str, doc_title: str = "StudyVault Document"):
    """Write standard multi-page PDF document using fpdf2."""
    if FPDF is None:
        raise ImportError("fpdf2 is not installed in the python environment.")
        
    pdf = StudyVaultPDF()
    pdf.doc_title = clean_pdf_text(doc_title)
    pdf.alias_nb_pages()
    pdf.add_page()
    pdf.set_font("helvetica", size=10)
    pdf.set_text_color(30, 30, 30)
    
    paragraphs = text.split('\n')
    for p in paragraphs:
        if p.strip():
            cleaned = clean_pdf_text(p)
            pdf.multi_cell(0, 6, cleaned)
            pdf.ln(3)
        else:
            pdf.ln(4)
            
    pdf.output(output_path)


def convert_to_pdf_direct(source_path: str, output_path: str, file_type: str) -> bool:
    """
    Directly convert formatted files (pptx, xlsx, docx) to PDF preserving page structures
    (e.g., slide-to-page mapping or sheet-to-page mapping).
    """
    file_type = file_type.lower()
    doc_title = os.path.basename(source_path)
    
    if file_type == 'pptx' and Presentation is not None and FPDF is not None:
        prs = Presentation(source_path)
        slide_width_emu = prs.slide_width
        slide_height_emu = prs.slide_height
        slide_width_mm = slide_width_emu / 36000.0
        slide_height_mm = slide_height_emu / 36000.0
        
        # FPDF orientation='L' swaps format, so pass height first then width
        pdf = StudyVaultPDF(orientation='L', unit='mm', format=(slide_height_mm, slide_width_mm))
        pdf.doc_title = clean_pdf_text(doc_title)
        pdf.is_slide = True
        pdf.set_auto_page_break(False)
        pdf.alias_nb_pages()
        
        for slide_idx, slide in enumerate(prs.slides):
            pdf.add_page()
            
            # Slide background fill if solid
            try:
                background = slide.background
                if background and background.fill and background.fill.type == 1:
                    bg_color = background.fill.fore_color
                    if hasattr(bg_color, "rgb") and bg_color.rgb is not None:
                        pdf.set_fill_color(bg_color.rgb[0], bg_color.rgb[1], bg_color.rgb[2])
                        pdf.rect(0, 0, slide_width_mm, slide_height_mm, 'F')
            except Exception:
                pass
            
            # Draw shapes
            for shape in slide.shapes:
                left = shape.left / 36000.0
                top = shape.top / 36000.0
                width = shape.width / 36000.0
                height = shape.height / 36000.0
                
                # Shape solid fill
                try:
                    if hasattr(shape, "fill") and shape.fill and shape.fill.type == 1:
                        fill_color = shape.fill.fore_color
                        if hasattr(fill_color, "rgb") and fill_color.rgb is not None:
                            pdf.set_fill_color(fill_color.rgb[0], fill_color.rgb[1], fill_color.rgb[2])
                            pdf.rect(left, top, width, height, 'F')
                except Exception:
                    pass
                
                # Render Pictures
                if hasattr(shape, "image") and shape.image is not None:
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        ext = image.ext or 'png'
                        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}") as tmp:
                            tmp.write(image_bytes)
                            tmp_name = tmp.name
                        pdf.image(tmp_name, x=left, y=top, w=width, h=height)
                        try:
                            os.unlink(tmp_name)
                        except OSError:
                            pass
                    except Exception as e:
                        print(f"Error rendering image shape: {e}", file=sys.stderr)
                
                # Render Tables
                elif hasattr(shape, "has_table") and shape.has_table:
                    try:
                        table = shape.table
                        row_top = top
                        for r_idx, row in enumerate(table.rows):
                            row_h = row.height / 36000.0
                            col_left = left
                            for c_idx, cell in enumerate(row.cells):
                                col_w = table.columns[c_idx].width / 36000.0
                                
                                # Draw cell bounding box
                                pdf.set_fill_color(248, 249, 250)
                                pdf.set_draw_color(180, 180, 180)
                                pdf.set_line_width(0.15)
                                pdf.rect(col_left, row_top, col_w, row_h, 'DF')
                                
                                if cell.text.strip():
                                    pdf.set_xy(col_left + 0.5, row_top + 0.5)
                                    pdf.set_font("helvetica", size=7)
                                    pdf.set_text_color(33, 37, 41)
                                    pdf.multi_cell(col_w - 1, 3, clean_pdf_text(cell.text))
                                col_left += col_w
                            row_top += row_h
                    except Exception as e:
                        print(f"Error rendering table shape: {e}", file=sys.stderr)
                
                # Render text boxes
                elif hasattr(shape, "text_frame") and shape.text.strip():
                    try:
                        is_title = 'title' in shape.name.lower() or 'header' in shape.name.lower() or slide_idx == 0 and 'title' in shape.name.lower()
                        if is_title:
                            pdf.set_font("helvetica", "B", 14)
                            pdf.set_text_color(20, 20, 20)
                            line_h = 7
                        else:
                            pdf.set_font("helvetica", size=9)
                            pdf.set_text_color(40, 40, 40)
                            line_h = 4.5
                        pdf.set_xy(left + 0.5, top + 0.5)
                        pdf.multi_cell(width - 1, line_h, clean_pdf_text(shape.text))
                    except Exception as e:
                        print(f"Error rendering text shape: {e}", file=sys.stderr)
        pdf.output(output_path)
        return True
        
    elif file_type == 'xlsx' and openpyxl is not None and FPDF is not None:
        wb = openpyxl.load_workbook(source_path, data_only=True, read_only=True)
        pdf = StudyVaultPDF()
        pdf.doc_title = clean_pdf_text(doc_title)
        pdf.alias_nb_pages()
        
        for sheet_name in wb.sheetnames:
            pdf.add_page()
            pdf.set_font("helvetica", "B", 14)
            pdf.cell(0, 10, f"Sheet: {sheet_name}", border=0, ln=1)
            pdf.ln(5)
            pdf.set_font("helvetica", size=9)
            
            sheet = wb[sheet_name]
            for row in sheet.iter_rows(values_only=True):
                row_vals = [str(cell).strip() for cell in row if cell is not None]
                if row_vals:
                    row_str = " | ".join(row_vals)
                    cleaned_row = clean_pdf_text(row_str)
                    pdf.multi_cell(0, 5, cleaned_row)
                    pdf.ln(1)
        pdf.output(output_path)
        return True
        
    elif file_type == 'docx' and docx is not None and FPDF is not None:
        # Simple docx parser mapping to PDF paragraphs
        doc = docx.Document(source_path)
        pdf = StudyVaultPDF()
        pdf.doc_title = clean_pdf_text(doc_title)
        pdf.alias_nb_pages()
        pdf.add_page()
        pdf.set_font("helvetica", size=10)
        
        for p in doc.paragraphs:
            if p.text.strip():
                cleaned = clean_pdf_text(p.text)
                pdf.multi_cell(0, 6, cleaned)
                pdf.ln(3)
                
        for table in doc.tables:
            pdf.ln(4)
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    row_str = " | ".join(row_text)
                    pdf.multi_cell(0, 5, clean_pdf_text(row_str))
                    pdf.ln(1)
        pdf.output(output_path)
        return True
        
    return False


def get_cached_pdf_path(source_path: str, file_type: str) -> str:
    """
    Get or create a cached PDF file path for non-PDF documents to support
    rich page layout viewing.
    """
    file_type = file_type.lower()
    if file_type == 'pdf':
        return source_path
        
    # Generate stable temp path based on file path hash
    path_hash = hashlib.md5(source_path.encode('utf-8')).hexdigest()
    temp_dir = tempfile.gettempdir()
    cached_pdf = os.path.join(temp_dir, f"studyvault_cache_{path_hash}.pdf")
    
    # Generate on-demand if missing or source is newer
    if not os.path.exists(cached_pdf) or os.path.getmtime(source_path) > os.path.getmtime(cached_pdf):
        # Try direct formatting conversion
        converted = False
        try:
            converted = convert_to_pdf_direct(source_path, cached_pdf, file_type)
        except Exception as e:
            print(f"Direct formatting PDF conversion failed: {e}", file=sys.stderr)
            
        if not converted:
            # Fallback to plain text extraction and standard PDF writing
            import extractor
            text, _ = extractor.extract_text_and_hash(source_path, file_type)
            clean_txt = clean_extracted_text(text)
            write_pdf(clean_txt, cached_pdf, os.path.basename(source_path))
            
    return cached_pdf
